"""Extract text + notes from each slide of the PPT for narration scripting."""
import json
from pathlib import Path
from pptx import Presentation

PPT = r"C:\Users\User\Downloads\DIGITALIZATON_4.pptx"
OUT = Path(r"C:\Users\User\Desktop\Collection_Auto\video_pipeline\slides_content.json")

prs = Presentation(PPT)
slides = []
for i, slide in enumerate(prs.slides, start=1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                if line:
                    texts.append(line)
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    slides.append({"index": i, "texts": texts, "notes": notes})

OUT.write_text(json.dumps(slides, indent=2, ensure_ascii=False), encoding="utf-8")
print(f"Slide count: {len(slides)}")
for s in slides:
    print(f"--- Slide {s['index']} ---")
    for t in s["texts"]:
        print(f"  • {t}")
    if s["notes"]:
        print(f"  NOTES: {s['notes'][:200]}")
