"""Check how many images each slide has."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r"C:\Users\User\Downloads\DIGITALIZATON_4.pptx")
for i, slide in enumerate(prs.slides, 1):
    if i > 35: break
    pics = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    print(f"Slide {i}: {pics} image(s)")
