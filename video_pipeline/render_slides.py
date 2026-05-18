"""Render PPT slides 1-35 to 1920x1080 PNGs by reading shapes via python-pptx
and laying them out with PIL. Preserves original layout positions and the
embedded equipment images, which is what matters for the narration video."""
import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

PPT = r"C:\Users\User\Downloads\DIGITALIZATON_4.pptx"
OUT_DIR = Path(r"C:\Users\User\Desktop\Collection_Auto\video_pipeline\slides")
OUT_DIR.mkdir(exist_ok=True, parents=True)

TARGET_W, TARGET_H = 1920, 1080
SLIDE_LIMIT = 35

BG = (245, 247, 250)
TITLE_BG = (12, 52, 102)
TITLE_FG = (255, 255, 255)
BODY_FG = (28, 35, 49)
ACCENT = (220, 38, 38)

FONT_CANDIDATES = [
    r"C:\Windows\Fonts\segoeuib.ttf",
    r"C:\Windows\Fonts\arialbd.ttf",
    r"C:\Windows\Fonts\seguisb.ttf",
]
REG_FONT_CANDIDATES = [
    r"C:\Windows\Fonts\segoeui.ttf",
    r"C:\Windows\Fonts\arial.ttf",
]

def load_font(paths, size):
    for p in paths:
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()

def emu_to_px(emu, slide_w_emu, slide_h_emu):
    return (
        int(emu[0] / slide_w_emu * TARGET_W),
        int(emu[1] / slide_h_emu * TARGET_H),
        int(emu[2] / slide_w_emu * TARGET_W),
        int(emu[3] / slide_h_emu * TARGET_H),
    )

def wrap_text(draw, text, font, max_width):
    words = text.split()
    if not words:
        return []
    lines = []
    cur = words[0]
    for w in words[1:]:
        test = cur + " " + w
        if draw.textlength(test, font=font) <= max_width:
            cur = test
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines

def render_slide(slide, idx, slide_w_emu, slide_h_emu):
    img = Image.new("RGB", (TARGET_W, TARGET_H), BG)
    draw = ImageDraw.Draw(img)

    # Header bar
    draw.rectangle([(0, 0), (TARGET_W, 90)], fill=TITLE_BG)
    header_font = load_font(FONT_CANDIDATES, 36)
    draw.text((40, 22), "Digitalization of Model Room — ZRTI Bhusawal",
              font=header_font, fill=TITLE_FG)
    # Slide number
    num_font = load_font(FONT_CANDIDATES, 28)
    draw.text((TARGET_W - 150, 28), f"Slide {idx}", font=num_font, fill=(180, 200, 230))

    # Footer bar
    draw.rectangle([(0, TARGET_H - 50), (TARGET_W, TARGET_H)], fill=TITLE_BG)

    # Collect text + image shapes
    text_blocks = []
    image_blocks = []
    for shape in slide.shapes:
        try:
            box = (shape.left, shape.top, shape.width, shape.height)
        except Exception:
            continue
        if any(b is None for b in box):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                image_blocks.append((box, blob))
            except Exception:
                pass
        elif shape.has_text_frame:
            texts = []
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    texts.append(t)
            if texts:
                text_blocks.append((box, texts))

    # Layout strategy: use original positions, scaled
    content_y0 = 110
    content_y1 = TARGET_H - 60

    # Render images at their original positions (scaled)
    for box, blob in image_blocks:
        x, y, w, h = emu_to_px(box, slide_w_emu, slide_h_emu)
        # Clamp to content area
        if y < content_y0: y = content_y0
        if y + h > content_y1: h = content_y1 - y
        if w <= 0 or h <= 0: continue
        try:
            pic = Image.open(io.BytesIO(blob)).convert("RGB")
            pic.thumbnail((w, h), Image.LANCZOS)
            # Center within the original box
            px = x + (w - pic.width) // 2
            py = y + (h - pic.height) // 2
            img.paste(pic, (px, py))
        except Exception as e:
            print(f"  image render error slide {idx}: {e}")

    # Render text blocks
    title_font = load_font(FONT_CANDIDATES, 56)
    body_font = load_font(REG_FONT_CANDIDATES, 38)
    small_body_font = load_font(REG_FONT_CANDIDATES, 30)

    for box, texts in text_blocks:
        x, y, w, h = emu_to_px(box, slide_w_emu, slide_h_emu)
        if y < content_y0: y = content_y0
        if w < 200: w = 200
        # Determine if this is a "title" block — heuristic: short text, top of slide
        avg_len = sum(len(t) for t in texts) / max(len(texts), 1)
        is_title = (len(texts) <= 2 and avg_len < 50 and y < TARGET_H * 0.4)

        font = title_font if is_title else body_font
        fill = ACCENT if is_title else BODY_FG

        cur_y = y
        for t in texts:
            # Use smaller font if doesn't fit
            f = font
            if not is_title and len(t) > 70:
                f = small_body_font
            wrapped = wrap_text(draw, t, f, w)
            for line in wrapped:
                if cur_y > content_y1 - 40:
                    break
                draw.text((x, cur_y), line, font=f, fill=fill)
                # Line height ~1.2x font size
                fs = f.size if hasattr(f, "size") else 36
                cur_y += int(fs * 1.25)
            cur_y += 8  # paragraph spacing

    # Footer text
    foot_font = load_font(REG_FONT_CANDIDATES, 22)
    draw.text((40, TARGET_H - 38), "Digitalization of Model Room • ZRTI Bhusawal",
              font=foot_font, fill=(180, 200, 230))

    out_path = OUT_DIR / f"slide_{idx:03d}.png"
    img.save(out_path, "PNG", optimize=True)
    return out_path

def main():
    prs = Presentation(PPT)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"Slide dim EMU: {slide_w} x {slide_h}")
    for i, slide in enumerate(prs.slides, start=1):
        if i > SLIDE_LIMIT:
            break
        path = render_slide(slide, i, slide_w, slide_h)
        print(f"Rendered {path.name}")

if __name__ == "__main__":
    main()
